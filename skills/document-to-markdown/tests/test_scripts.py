"""Tests for MarkItDown skill scripts."""

import subprocess
import sys
from pathlib import Path

import pytest

SCRIPTS = Path(__file__).parent.parent / "scripts"


def run(script: str, *args: str) -> subprocess.CompletedProcess:
    return subprocess.run(
        [sys.executable, str(SCRIPTS / script), *args],
        capture_output=True,
        text=True,
    )


@pytest.fixture(scope="session")
def sample_docx(tmp_path_factory) -> Path:
    from docx import Document

    path = tmp_path_factory.mktemp("docx") / "sample.docx"
    document = Document()
    document.add_heading("Quarterly Update", level=1)
    document.add_paragraph("The launch is on track.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Projects"
    table.cell(1, 1).text = "12"
    document.save(path)
    return path


@pytest.fixture(scope="session")
def sample_pdf(tmp_path_factory) -> Path:
    from reportlab.pdfgen.canvas import Canvas

    path = tmp_path_factory.mktemp("pdf") / "sample.pdf"
    canvas = Canvas(str(path))
    canvas.drawString(72, 720, "Quarterly Update PDF")
    canvas.save()
    return path


@pytest.fixture(scope="session")
def sample_pptx(tmp_path_factory) -> Path:
    from pptx import Presentation

    path = tmp_path_factory.mktemp("pptx") / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Quarterly Update Slides"
    presentation.save(path)
    return path


@pytest.fixture(scope="session")
def sample_xlsx(tmp_path_factory) -> Path:
    from openpyxl import Workbook

    path = tmp_path_factory.mktemp("xlsx") / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["Metric", "Value"])
    sheet.append(["Projects", 12])
    workbook.save(path)
    return path


class TestConvert:
    def test_converts_text_file(self, tmp_path):
        source = tmp_path / "notes.txt"
        source.write_text("A short note for analysis.")
        output = tmp_path / "nested" / "notes.md"

        result = run("convert.py", str(source), "--output", str(output))

        assert result.returncode == 0
        assert output.read_text() == "A short note for analysis."
        assert "Converted:" in result.stdout

    @pytest.mark.parametrize(
        ("fixture_name", "expected_text"),
        [
            ("sample_docx", "Quarterly Update"),
            ("sample_pdf", "Quarterly Update PDF"),
            ("sample_pptx", "Quarterly Update Slides"),
            ("sample_xlsx", "Projects"),
        ],
    )
    def test_converts_office_and_pdf_content(
        self,
        request,
        fixture_name,
        expected_text,
        tmp_path,
    ):
        source = request.getfixturevalue(fixture_name)
        output = tmp_path / "update.md"

        result = run("convert.py", str(source), "--output", str(output))

        assert result.returncode == 0
        content = output.read_text()
        assert expected_text in content

    def test_rejects_urls_as_non_local_inputs(self, tmp_path):
        result = run(
            "convert.py",
            "https://example.com/document.pdf",
            "--output",
            str(tmp_path / "document.md"),
        )

        assert result.returncode != 0
        assert "Input file not found" in result.stderr

    def test_rejects_unsupported_file_types(self, tmp_path):
        source = tmp_path / "archive.zip"
        source.write_bytes(b"not an archive")

        result = run("convert.py", str(source), "--output", str(tmp_path / "archive.md"))

        assert result.returncode != 0
        assert "Unsupported input type" in result.stderr

    def test_does_not_overwrite_existing_output_without_force(self, tmp_path):
        source = tmp_path / "notes.txt"
        source.write_text("New note")
        output = tmp_path / "notes.md"
        output.write_text("Existing note")

        result = run("convert.py", str(source), "--output", str(output))

        assert result.returncode != 0
        assert output.read_text() == "Existing note"

    def test_force_overwrites_existing_output(self, tmp_path):
        source = tmp_path / "notes.txt"
        source.write_text("New note")
        output = tmp_path / "notes.md"
        output.write_text("Existing note")

        result = run("convert.py", str(source), "--output", str(output), "--force")

        assert result.returncode == 0
        assert output.read_text() == "New note"
