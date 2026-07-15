"""Tests for powerpoint skill Python scripts."""

import json
import subprocess
import sys
from pathlib import Path

import pytest

SCRIPTS = Path(__file__).parent.parent / "scripts"


def run(script: str, *args: str) -> subprocess.CompletedProcess:
    return subprocess.run(
        [sys.executable, str(SCRIPTS / script), *args],
        capture_output=True,
        text=True,
    )


@pytest.fixture(scope="session")
def sample_pptx(tmp_path_factory) -> Path:
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path_factory.mktemp("pptx") / "sample.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Presentation"
    slide.placeholders[1].text = "Subtitle text"

    layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout2)
    slide2.shapes.title.text = "Slide Two"
    slide2.placeholders[1].text = "Content here"

    prs.save(str(path))
    return path


class TestExtractText:
    def test_extracts_slide_text(self, sample_pptx):
        result = run("extract-text.py", str(sample_pptx))
        assert result.returncode == 0
        assert "Test Presentation" in result.stdout
        assert "Slide Two" in result.stdout

    def test_json_output(self, sample_pptx):
        result = run("extract-text.py", str(sample_pptx), "--json")
        assert result.returncode == 0
        data = json.loads(result.stdout)
        assert data["slideCount"] == 2
        assert any("Test Presentation" in t for t in data["slides"][0]["texts"])


class TestCreate:
    def test_creates_from_json_string(self, tmp_path):
        pytest.importorskip("pptx")
        out = tmp_path / "deck.pptx"
        slides = json.dumps([
            {"layout": "title", "title": "Hello", "content": "World"},
            {"layout": "content", "title": "Points", "content": "- First\n- Second"},
        ])
        result = run("create.py", "--content", slides, "--output", str(out))
        assert result.returncode == 0
        assert out.exists()
        from pptx import Presentation
        prs = Presentation(str(out))
        assert len(prs.slides) == 2

    def test_creates_from_file(self, tmp_path):
        pytest.importorskip("pptx")
        slides_file = tmp_path / "slides.json"
        slides_file.write_text(json.dumps([
            {"title": "Slide 1", "content": "Body"},
        ]))
        out = tmp_path / "from_file.pptx"
        result = run("create.py", "--input", str(slides_file), "--output", str(out))
        assert result.returncode == 0
        assert out.exists()

    def test_rejects_unsupported_layout(self, tmp_path):
        out = tmp_path / "unsupported-layout.pptx"
        slides = json.dumps([{"layout": "two-col", "title": "Unsupported"}])

        result = run("create.py", "--content", slides, "--output", str(out))

        assert result.returncode != 0
        assert "Unsupported layout" in result.stderr


class TestInventory:
    def test_outputs_shape_map(self, sample_pptx):
        result = run("inventory.py", str(sample_pptx))
        assert result.returncode == 0
        data = json.loads(result.stdout)
        assert "slide-0" in data
        assert any(
            any("Test Presentation" in p["text"] for p in shapes.get("paragraphs", []))
            for shapes in data["slide-0"].values()
        )

    def test_keeps_duplicate_shape_names_addressable(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches

        source = tmp_path / "duplicate-shapes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for index, text in enumerate(("First", "Second")):
            shape = slide.shapes.add_textbox(Inches(index), Inches(1), Inches(1), Inches(1))
            shape.name = "Duplicate shape"
            shape.text = text
        prs.save(source)

        result = run("inventory.py", str(source))

        assert result.returncode == 0
        shapes = json.loads(result.stdout)["slide-0"]
        duplicates = [
            (key, value)
            for key, value in shapes.items()
            if value["name"] == "Duplicate shape"
        ]
        assert [key for key, _ in duplicates] == ["shape-0", "shape-1"]
        assert [value["paragraphs"][0]["text"] for _, value in duplicates] == ["First", "Second"]


class TestReplace:
    def test_find_replace(self, sample_pptx, tmp_path):
        pytest.importorskip("pptx")
        out = tmp_path / "replaced.pptx"
        result = run("replace.py", str(sample_pptx), str(out), "--find", "Slide Two", "--replace", "Slide 2")
        assert result.returncode == 0
        assert out.exists()
        from pptx import Presentation
        prs = Presentation(str(out))
        all_text = " ".join(
            "".join(r.text for r in para.runs)
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
        assert "Slide 2" in all_text

    def test_find_replace_across_runs(self, tmp_path):
        from pptx import Presentation

        source = tmp_path / "split-runs.pptx"
        out = tmp_path / "replaced.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_frame = slide.shapes.add_textbox(0, 0, 100, 100).text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.add_run().text = "Draft "
        paragraph.add_run().text = "Version"
        prs.save(source)

        result = run(
            "replace.py",
            str(source),
            str(out),
            "--find",
            "Draft Version",
            "--replace",
            "Final Version",
        )

        assert result.returncode == 0
        assert "Replaced 1 occurrence" in result.stdout
        replaced = Presentation(out)
        assert replaced.slides[0].shapes[0].text == "Final Version"

    def test_find_replace_in_table_cells(self, tmp_path):
        from pptx import Presentation

        source = tmp_path / "table.pptx"
        out = tmp_path / "replaced.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(1, 1, 0, 0, 100, 100).table
        table.cell(0, 0).text = "Draft table"
        presentation.save(source)

        result = run(
            "replace.py",
            str(source),
            str(out),
            "--find",
            "Draft",
            "--replace",
            "Final",
        )

        assert result.returncode == 0
        assert "Replaced 1 occurrence" in result.stdout
        replaced = Presentation(out)
        assert replaced.slides[0].shapes[0].table.cell(0, 0).text == "Final table"

    def test_inventory_replace_removes_omitted_paragraphs(self, sample_pptx, tmp_path):
        from pptx import Presentation

        inventory = json.loads(run("inventory.py", str(sample_pptx)).stdout)
        shape_key, _ = next(
            (key, value)
            for key, value in inventory["slide-1"].items()
            if any(paragraph["text"] == "Slide Two" for paragraph in value["paragraphs"])
        )
        replacements = {"slide-1": {shape_key: {"paragraphs": [{"text": "Updated", "level": 0}]}}}
        replacements_path = tmp_path / "replacements.json"
        replacements_path.write_text(json.dumps(replacements))
        out = tmp_path / "replaced-inventory.pptx"

        result = run("replace.py", str(sample_pptx), str(replacements_path), str(out))

        assert result.returncode == 0
        all_text = "\n".join(
            shape.text
            for shape in Presentation(out).slides[1].shapes
            if shape.has_text_frame
        )
        assert "Updated" in all_text
        assert "Slide Two" not in all_text
