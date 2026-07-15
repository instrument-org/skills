#!/usr/bin/env python3
"""Inventory all text shapes in a .pptx file.

Outputs a JSON map of slide -> shape -> text content and position.
Use this before replace.py to understand what text can be substituted.

Usage:
  python scripts/inventory.py input.pptx [--output inventory.json]
"""

import argparse
import json
import sys
from pathlib import Path


def main():
    parser = argparse.ArgumentParser(description="Inventory text shapes in a .pptx")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("--output", help="Save JSON to file (default: stdout)")
    args = parser.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        sys.exit("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(args.input)
    result = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_key = f"slide-{slide_idx}"
        shapes = {}
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs)
                if text.strip():
                    paragraphs.append({
                        "text": text,
                        "level": para.level,
                        "bold": any(r.font.bold for r in para.runs if r.font.bold is not None),
                    })
            if not paragraphs:
                continue
            shape_key = f"shape-{shape_idx}"
            shapes[shape_key] = {
                "name": shape.name,
                "left": shape.left / 914400 if shape.left else 0,
                "top": shape.top / 914400 if shape.top else 0,
                "width": shape.width / 914400 if shape.width else 0,
                "height": shape.height / 914400 if shape.height else 0,
                "paragraphs": paragraphs,
            }
        if shapes:
            result[slide_key] = shapes

    output = json.dumps(result, indent=2)
    if args.output:
        Path(args.output).write_text(output)
        print(f"Saved: {args.output}")
    else:
        print(output)


if __name__ == "__main__":
    main()
