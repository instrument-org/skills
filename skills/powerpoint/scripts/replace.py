#!/usr/bin/env python3
"""Replace text in a .pptx presentation using an inventory JSON.

Workflow:
  1. Run inventory.py to get the shape map
  2. Edit the JSON to set new text for each shape
  3. Run replace.py to apply changes

The replacements JSON has the same structure as inventory.py output.
Only shapes listed under 'paragraphs' are modified; other shapes are untouched.

Usage:
  python scripts/replace.py input.pptx replacements.json output.pptx

Or for simple find-and-replace across all text:
  python scripts/replace.py input.pptx --find "Draft" --replace "Final" output.pptx
"""

import argparse
import json
import sys


def replace_via_inventory(prs, replacements: dict):
    """Apply inventory-style replacements."""
    for slide_idx, slide in enumerate(prs.slides):
        slide_key = f"slide-{slide_idx}"
        if slide_key not in replacements:
            continue
        shape_replacements = replacements[slide_key]
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            data = shape_replacements.get(f"shape-{shape_idx}")
            if data is None:
                continue
            new_paras = data.get("paragraphs", [])
            tf = shape.text_frame
            tf.clear()
            for i, para_data in enumerate(new_paras):
                if i == 0:
                    para = tf.paragraphs[i]
                else:
                    para = tf.add_paragraph()
                if para.runs:
                    para.runs[0].text = para_data.get("text", "")
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    run = para.add_run()
                    run.text = para_data.get("text", "")
                para.level = para_data.get("level", 0)


def replace_text(paragraph, find: str, replacement: str):
    """Replace text in a paragraph, including matches spanning formatting runs."""
    count = 0
    search_from = 0

    while True:
        runs = list(paragraph.runs)
        text = "".join(run.text for run in runs)
        start = text.find(find, search_from)
        if start == -1:
            return count

        end = start + len(find)
        offset = 0
        start_index = None
        end_index = None
        start_offset = 0
        end_offset = 0

        for index, run in enumerate(runs):
            run_end = offset + len(run.text)
            if start_index is None and start < run_end:
                start_index = index
                start_offset = start - offset
            if end <= run_end:
                end_index = index
                end_offset = end - offset
                break
            offset = run_end

        if start_index is None or end_index is None:
            return count

        start_run = runs[start_index]
        end_run = runs[end_index]
        if start_run == end_run:
            start_run.text = (
                start_run.text[:start_offset]
                + replacement
                + start_run.text[end_offset:]
            )
        else:
            start_run.text = start_run.text[:start_offset] + replacement
            for run in runs[start_index + 1:end_index]:
                run.text = ""
            end_run.text = end_run.text[end_offset:]

        count += 1
        search_from = start + len(replacement)


def replace_find_replace(prs, find: str, replace: str):
    """Simple find-and-replace across all text shapes."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    count += replace_text(para, find, replace)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            count += replace_text(para, find, replace)
    return count


def main():
    parser = argparse.ArgumentParser(description="Replace text in a .pptx file")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("replacements_or_output",
                        help="Replacements JSON file (inventory mode) or output path (find/replace mode)")
    parser.add_argument("output", nargs="?", help="Output .pptx file (inventory mode)")
    parser.add_argument("--find", help="Text to find (simple mode)")
    parser.add_argument("--replace", dest="replace_with", help="Replacement text (simple mode)")
    args = parser.parse_args()
    if (args.find is None) != (args.replace_with is None):
        parser.error("--find and --replace must be supplied together")
    if args.find == "":
        parser.error("--find must not be empty")

    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(args.input)

    if args.find is not None:
        count = replace_find_replace(prs, args.find, args.replace_with)
        output_path = args.replacements_or_output
        prs.save(output_path)
        print(f"Replaced {count} occurrence(s) -> {output_path}")
    else:
        with open(args.replacements_or_output) as f:
            replacements = json.load(f)
        output_path = args.output
        if not output_path:
            sys.exit("Provide output path as third argument in inventory mode")
        replace_via_inventory(prs, replacements)
        prs.save(output_path)
        print(f"Applied replacements -> {output_path}")


if __name__ == "__main__":
    main()
