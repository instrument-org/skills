#!/usr/bin/env python3
"""Extract text from a PowerPoint presentation (.pptx)."""

import argparse
import json
import sys


def main():
    parser = argparse.ArgumentParser(description="Extract text from a .pptx file")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("--json", action="store_true", dest="as_json",
                        help="Output structured JSON with per-slide text")
    args = parser.parse_args()

    try:
        from pptx import Presentation
    except ImportError:
        sys.exit(
            "python-pptx is unavailable. Reload this skill to retry dependency setup."
        )

    prs = Presentation(args.input)

    if args.as_json:
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in para.runs).strip()
                        if text:
                            texts.append(text)
            slides.append({"slide": i, "texts": texts})
        print(json.dumps({"slideCount": len(prs.slides), "slides": slides}, indent=2))
    else:
        for i, slide in enumerate(prs.slides, start=1):
            print(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in para.runs).strip()
                        if text:
                            print(text)


if __name__ == "__main__":
    main()
