#!/usr/bin/env python3
"""Create a PowerPoint presentation (.pptx) from a JSON slide definition.

JSON format — array of slide objects:
  [
    {
      "layout": "title",       // "title", "content", or "blank" (default: "content")
      "title": "Slide Title",
      "content": "Body text or bullet points separated by newlines",
      "notes": "Speaker notes (optional)"
    }
  ]

Bullet points: any content line starting with "- " or "• " becomes a bullet.

Examples:
  python scripts/create.py --input slides.json --output deck.pptx
  python scripts/create.py --content '[{"title":"Hello","content":"World"}]' --output hello.pptx
"""

import argparse
import json
import sys


SLIDE_WIDTH_EMU = 9144000   # 10 inches
SLIDE_HEIGHT_EMU = 5143500  # 5.625 inches


def add_slide(prs, slide_def: dict):
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    layout_name = slide_def.get("layout", "content")
    title_text = slide_def.get("title", "")
    content_text = slide_def.get("content", "")
    notes_text = slide_def.get("notes", "")

    layout_map = {"title": 0, "content": 1, "blank": 6}
    if layout_name not in layout_map:
        sys.exit(
            f"Unsupported layout: {layout_name}. "
            f"Use one of: {', '.join(layout_map)}"
        )
    layout_idx = layout_map[layout_name]
    layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Set title
    if title_text and slide.shapes.title:
        slide.shapes.title.text = title_text

    # Set body content
    body_placeholder = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_placeholder = ph
            break

    if body_placeholder and content_text:
        tf = body_placeholder.text_frame
        tf.clear()
        lines = content_text.splitlines()
        first = True
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            if stripped.startswith(("- ", "• ")):
                para.text = stripped.lstrip("-• ").strip()
                para.level = 1
            else:
                para.text = stripped
                para.level = 0

    # Speaker notes
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    return slide


def main():
    parser = argparse.ArgumentParser(description="Create a PowerPoint presentation")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    parser.add_argument("--input", help="JSON file with slide definitions")
    parser.add_argument("--content", help="JSON string with slide definitions")
    parser.add_argument("--title", help="Presentation title (metadata)")
    args = parser.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        sys.exit(
            "python-pptx is unavailable. Reload this skill to retry dependency setup."
        )

    if args.input:
        with open(args.input) as f:
            slides_def = json.load(f)
    elif args.content:
        slides_def = json.loads(args.content)
    else:
        sys.exit("Provide --input or --content")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    for slide_def in slides_def:
        add_slide(prs, slide_def)

    if args.title:
        prs.core_properties.title = args.title

    prs.save(args.output)
    print(f"Created: {args.output} ({len(slides_def)} slides)")


if __name__ == "__main__":
    main()
